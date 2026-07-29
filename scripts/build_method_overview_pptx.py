#!/usr/bin/env python3
"""Build an editable PowerPoint version of the paper's method overview."""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "paper_aaai2027" / "figures" / "method_overview_editable.pptx"

SLIDE_W_IN = 13.5
SCALE = SLIDE_W_IN / 18.0
SLIDE_H_IN = 10.7 * SCALE

INK = "202124"
WHITE = "FFFFFF"
SOFTGRAY = "E7E8EA"
LIGHTGRAY = "F7F8F9"
MIDGRAY = "9AA0A6"
BLUE = "2A6FBB"
BLUEFILL = "DCEBFA"
PURPLE = "7251B5"
PURPLEFILL = "E9DDF7"
ORANGE = "F06A14"
ORANGEFILL = "FDE3CC"
GREEN = "3F7D2A"
GREENFILL = "E4F0DD"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def x(value: float):
    return Inches(value * SCALE)


def y_from_top(value: float):
    return Inches(value * SCALE)


def top_from_y(y_top: float):
    return Inches((10.7 - y_top) * SCALE)


def set_name(shape, name: str) -> None:
    element = shape._element
    if hasattr(element, "nvSpPr"):
        element.nvSpPr.cNvPr.set("name", name)
    elif hasattr(element, "nvGrpSpPr"):
        element.nvGrpSpPr.cNvPr.set("name", name)


def set_arrowhead(connector, *, head: bool = True, tail: bool = False) -> None:
    ln = connector._element.spPr.get_or_add_ln()
    for tag, enabled in (("a:headEnd", head), ("a:tailEnd", tail)):
        node = ln.find(tag, ln.nsmap)
        if node is None:
            node = OxmlElement(tag)
            ln.append(node)
        node.set("type", "triangle" if enabled else "none")
        node.set("w", "sm")
        node.set("len", "sm")


class Canvas:
    def __init__(self, slide):
        self.slide = slide
        self.current: list = []

    def track(self, shape, name: str | None = None):
        if name:
            set_name(shape, name)
        self.current.append(shape)
        return shape

    def clear_group(self) -> None:
        self.current = []

    def group(self, name: str):
        grp = self.slide.shapes.add_group_shape(self.current)
        set_name(grp, name)
        self.current = []
        return grp

    def rect(
        self,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        *,
        fill: str = WHITE,
        line: str = MIDGRAY,
        width: float = 0.5,
        rounded: bool = False,
        dashed: bool = False,
        name: str | None = None,
    ):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
        shp = self.slide.shapes.add_shape(
            shape_type,
            x(x1),
            top_from_y(y2),
            x(x2 - x1),
            x(y2 - y1),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(width)
        if dashed:
            shp.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        return self.track(shp, name)

    def circle(
        self,
        cx: float,
        cy: float,
        diameter: float,
        *,
        fill: str = WHITE,
        line: str = INK,
        width: float = 0.5,
        name: str | None = None,
    ):
        shp = self.slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x(cx - diameter / 2),
            top_from_y(cy + diameter / 2),
            x(diameter),
            x(diameter),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(width)
        return self.track(shp, name)

    def text(
        self,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        value: str,
        *,
        size: float = 8.0,
        color: str = INK,
        bold: bool = False,
        align: PP_ALIGN = PP_ALIGN.CENTER,
        valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
        font: str = "Arial",
        margins: float = 0.01,
        name: str | None = None,
    ):
        shp = self.slide.shapes.add_textbox(
            x(x1),
            top_from_y(y2),
            x(x2 - x1),
            x(y2 - y1),
        )
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = x(margins)
        tf.margin_right = x(margins)
        tf.margin_top = x(margins)
        tf.margin_bottom = x(margins)
        tf.vertical_anchor = valign
        lines = value.split("\n")
        for idx, line_value in enumerate(lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = line_value
            p.alignment = align
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.0
            for run in p.runs:
                run.font.name = font
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = rgb(color)
        return self.track(shp, name)

    def labeled_rect(
        self,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        value: str,
        *,
        fill: str = WHITE,
        line: str = MIDGRAY,
        width: float = 0.5,
        rounded: bool = True,
        size: float = 8.0,
        color: str = INK,
        bold: bool = False,
        name: str | None = None,
    ):
        shp = self.rect(
            x1,
            y1,
            x2,
            y2,
            fill=fill,
            line=line,
            width=width,
            rounded=rounded,
            name=name,
        )
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = x(0.04)
        tf.margin_right = x(0.04)
        tf.margin_top = x(0.02)
        tf.margin_bottom = x(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for idx, line_value in enumerate(value.split("\n")):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = line_value
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.0
            for run in p.runs:
                run.font.name = "Arial"
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.color.rgb = rgb(color)
        return shp

    def line(
        self,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        *,
        color: str = INK,
        width: float = 0.5,
        dashed: bool = False,
        arrow: bool = False,
        name: str | None = None,
    ):
        conn = self.slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            x(x1),
            top_from_y(y1),
            x(x2),
            top_from_y(y2),
        )
        conn.line.color.rgb = rgb(color)
        conn.line.width = Pt(width)
        if dashed:
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
        if arrow:
            set_arrowhead(conn, head=True)
        return self.track(conn, name)

    def arrow_shape(
        self,
        x1: float,
        y1: float,
        x2: float,
        y2: float,
        *,
        direction: str = "right",
        fill: str = INK,
        name: str | None = None,
    ):
        kind = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
        shp = self.slide.shapes.add_shape(
            kind,
            x(x1),
            top_from_y(y2),
            x(x2 - x1),
            x(y2 - y1),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
        shp.line.fill.background()
        return self.track(shp, name)


def add_step_header(c: Canvas, step: int, cx: float, title: str, title_box: tuple[float, float, float, float]):
    c.circle(cx, 9.46, 0.39, fill=INK, line=INK, name=f"Stage {step} number")
    c.text(
        cx - 0.14,
        9.32,
        cx + 0.14,
        9.60,
        str(step),
        size=13.5,
        color=WHITE,
        bold=True,
        name=f"Stage {step} number text",
    )
    c.text(*title_box, title, size=14.0, bold=True, name=f"Stage {step} title")


def add_token_cell(
    c: Canvas,
    x0: float,
    y0: float,
    w: float,
    h: float,
    *,
    fill: str,
    line: str = MIDGRAY,
    crossed: bool = False,
    name: str,
):
    c.rect(x0, y0, x0 + w, y0 + h, fill=fill, line=line, width=0.35, name=name)
    if crossed:
        c.line(x0 + 0.03, y0 + 0.03, x0 + w - 0.03, y0 + h - 0.03, color=MIDGRAY, width=0.3)
        c.line(x0 + w - 0.03, y0 + 0.03, x0 + 0.03, y0 + h - 0.03, color=MIDGRAY, width=0.3)


def build_stage_1(c: Canvas) -> None:
    c.clear_group()
    c.rect(0.18, 2.18, 3.55, 9.82, fill=WHITE, line=INK, width=1.0, rounded=True, name="Stage 1 panel")
    add_step_header(c, 1, 0.62, "OCR-Critical Query", (0.88, 9.26, 3.30, 9.66))

    c.rect(0.45, 4.28, 3.28, 8.96, fill=LIGHTGRAY, line=MIDGRAY, width=0.65, rounded=True, name="Editable invoice")
    c.text(0.62, 8.48, 1.70, 8.82, "INVOICE", size=14.0, bold=True, align=PP_ALIGN.LEFT, name="Invoice title")
    c.rect(0.59, 8.44, 1.55, 8.83, fill=WHITE, line=ORANGE, width=0.9, name="Invoice title bbox")
    c.text(2.04, 8.47, 3.08, 8.75, "2024-05-14", size=9.0, align=PP_ALIGN.RIGHT, name="Invoice date")
    c.rect(2.09, 8.43, 3.12, 8.81, fill=WHITE, line=ORANGE, width=0.9, name="Invoice date bbox")
    c.text(0.63, 7.82, 1.95, 8.18, "ID  INV-98231", size=10.5, bold=True, align=PP_ALIGN.LEFT)
    c.rect(0.57, 7.76, 1.92, 8.18, fill=WHITE, line=ORANGE, width=0.9, name="Invoice number bbox")
    c.text(0.63, 7.18, 2.34, 7.56, "Bill to: Acme Corp.", size=10.5, align=PP_ALIGN.LEFT, name="Invoice address")
    c.rect(0.57, 7.10, 2.42, 7.62, fill=WHITE, line=ORANGE, width=0.9, name="Invoice address bbox")

    c.rect(0.61, 5.33, 3.12, 6.30, fill=WHITE, line=MIDGRAY, width=0.45, name="Invoice item table")
    c.line(0.61, 5.96, 3.12, 5.96, color=MIDGRAY, width=0.4)
    c.line(0.61, 5.64, 3.12, 5.64, color=MIDGRAY, width=0.4)
    c.line(1.86, 5.33, 1.86, 6.30, color=MIDGRAY, width=0.4)
    c.line(2.37, 5.33, 2.37, 6.30, color=MIDGRAY, width=0.4)
    c.text(0.72, 5.68, 1.68, 6.02, "items", size=9.5, color=MIDGRAY, align=PP_ALIGN.LEFT)
    c.text(2.34, 5.68, 3.08, 6.02, "$35.97", size=9.5, color=MIDGRAY, align=PP_ALIGN.RIGHT)
    c.text(1.48, 4.86, 2.28, 5.18, "Total", size=10.5, bold=True, align=PP_ALIGN.LEFT)
    c.text(2.22, 4.84, 3.06, 5.20, "$41.36", size=11.5, bold=True, align=PP_ALIGN.RIGHT)
    c.rect(1.48, 4.84, 3.10, 5.19, fill=WHITE, line=ORANGE, width=0.9, name="Invoice total bbox")
    c.text(0.72, 4.36, 3.00, 4.70, "localized OCR evidence", size=10.5, color=ORANGE, bold=True)

    c.circle(0.65, 3.24, 0.26, fill=WHITE, line=INK, width=0.8, name="User head")
    c.line(0.42, 2.86, 0.88, 2.86, color=INK, width=0.8)
    c.line(0.42, 2.86, 0.48, 3.02, color=INK, width=0.8)
    c.line(0.88, 2.86, 0.82, 3.02, color=INK, width=0.8)
    c.labeled_rect(
        1.03,
        2.67,
        3.30,
        3.45,
        "Is the target text\npresent?",
        fill=LIGHTGRAY,
        line=BLUE,
        width=0.8,
        size=12.5,
        name="Editable query bubble",
    )
    c.group("Stage 1 - OCR-critical query")


def build_stage_2(c: Canvas) -> None:
    c.clear_group()
    c.rect(3.85, 2.18, 7.12, 9.82, fill=WHITE, line=INK, width=1.0, rounded=True, name="Stage 2 panel")
    add_step_header(c, 2, 4.28, "Projected\nVisual Tokens", (4.55, 9.12, 6.82, 9.78))
    c.text(4.00, 8.68, 7.00, 9.04, "Vision encoder + projector", size=12.0)

    encoder = c.slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, x(4.23), top_from_y(8.77), x(0.43), x(0.60))
    encoder.fill.solid()
    encoder.fill.fore_color.rgb = rgb(LIGHTGRAY)
    encoder.line.color.rgb = rgb(BLUE)
    encoder.line.width = Pt(1.0)
    c.track(encoder, "Editable vision encoder")
    c.arrow_shape(4.72, 8.38, 5.12, 8.56, fill=INK)
    for idx in range(5):
        c.rect(5.23 + 0.23 * idx, 8.34, 5.38 + 0.23 * idx, 8.62, fill=BLUEFILL, line=BLUE, width=0.55, name=f"Projected token {idx+1}")
    c.arrow_shape(6.38, 8.38, 6.75, 8.56, fill=INK)
    c.arrow_shape(5.39, 7.72, 5.59, 8.12, direction="down", fill=INK)

    evidence = {(1, 8), (2, 8), (1, 7), (1, 6), (1, 5), (5, 4), (6, 4), (7, 4), (8, 4), (2, 0), (3, 0), (4, 0), (5, 0), (6, 0), (7, 0)}
    for row in range(9):
        for col in range(9):
            is_evidence = (col, row) in evidence
            add_token_cell(
                c,
                4.12 + 0.31 * col,
                4.48 + 0.31 * row,
                0.23,
                0.23,
                fill=ORANGEFILL if is_evidence else SOFTGRAY,
                line=ORANGE if is_evidence else MIDGRAY,
                name=f"Projected grid r{row+1} c{col+1}",
            )
    c.rect(4.25, 3.73, 4.47, 3.95, fill=SOFTGRAY, line=MIDGRAY, width=0.45)
    c.text(4.60, 3.64, 6.70, 4.04, "ordinary token", size=11.5, align=PP_ALIGN.LEFT)
    c.rect(4.25, 3.30, 4.47, 3.52, fill=ORANGEFILL, line=ORANGE, width=0.45)
    c.text(4.60, 3.10, 6.96, 3.65, "OCR evidence token\n(overlaps text bbox)", size=11.0, align=PP_ALIGN.LEFT)
    c.group("Stage 2 - projected visual tokens")


def build_stage_3(c: Canvas) -> None:
    c.clear_group()
    c.rect(7.42, 2.18, 13.55, 9.82, fill=WHITE, line=INK, width=1.0, rounded=True, name="Stage 3 panel")
    add_step_header(c, 3, 7.86, "Budgeted Token Selection", (8.18, 9.25, 13.15, 9.66))

    signals = [
        ("Target-\nconditioned\nrelevance", BLUE, BLUEFILL, 8.16, 8.91),
        ("Projected-token\nnorm", PURPLE, PURPLEFILL, 7.32, 7.96),
        ("Evidence / grid\ncoverage prior", ORANGE, ORANGEFILL, 6.40, 7.04),
    ]
    for idx, (label, color, fill, y1, y2) in enumerate(signals):
        c.labeled_rect(7.72, y1, 9.50, y2, label, fill=fill, line=color, width=0.9, size=12.0, color=color, bold=True, name=f"Scoring signal {idx+1}")
        for row in range(3):
            for col in range(3):
                add_token_cell(
                    c,
                    9.82 + 0.22 * col,
                    y1 - 0.01 + 0.22 * row,
                    0.16,
                    0.16,
                    fill=fill,
                    line=color,
                    name=f"Signal {idx+1} grid r{row+1} c{col+1}",
                )
        c.line(9.43, (y1 + y2) / 2, 9.76, (y1 + y2) / 2, color=color, width=0.8, dashed=True, arrow=True)

    c.labeled_rect(
        10.66,
        6.89,
        13.20,
        8.36,
        "s_i = relevance\n+ token norm\n+ evidence prior",
        fill=WHITE,
        line=INK,
        width=0.9,
        size=13.0,
        name="Editable score formula",
    )
    c.line(10.50, 8.50, 10.66, 8.05, color=BLUE, width=0.8, arrow=True)
    c.line(10.50, 7.58, 10.66, 7.58, color=PURPLE, width=0.8, arrow=True)
    c.line(10.50, 6.66, 10.66, 7.13, color=ORANGE, width=0.8, arrow=True)
    c.arrow_shape(11.84, 6.36, 12.02, 6.86, direction="down", fill=INK)

    c.text(7.72, 5.74, 10.75, 6.14, "Budgeted top-k selection", size=12.0, bold=True)
    evidence_before = {(1, 6), (2, 6), (1, 5), (1, 4), (5, 4), (6, 4), (7, 4), (3, 1), (4, 1), (5, 1), (6, 1)}
    for row in range(7):
        for col in range(8):
            is_evidence = (col, row) in evidence_before
            add_token_cell(
                c,
                7.78 + 0.24 * col,
                3.87 + 0.24 * row,
                0.18,
                0.18,
                fill=ORANGEFILL if is_evidence else SOFTGRAY,
                line=ORANGE if is_evidence else MIDGRAY,
                name=f"Before pruning r{row+1} c{col+1}",
            )
    c.arrow_shape(9.84, 4.57, 10.34, 4.77, fill=INK, name="Pruning arrow")
    c.text(9.70, 4.82, 10.48, 5.25, "prune", size=10.5, bold=True)

    selected_evidence = {(1, 6), (1, 5), (1, 4), (6, 4), (7, 4), (3, 1), (4, 1), (5, 1), (6, 1)}
    selected_other = {(8, 6), (4, 5), (2, 3), (7, 2), (0, 0)}
    for row in range(7):
        for col in range(9):
            is_evidence = (col, row) in selected_evidence
            is_other = (col, row) in selected_other
            add_token_cell(
                c,
                10.50 + 0.23 * col,
                3.87 + 0.24 * row,
                0.17,
                0.18,
                fill=ORANGEFILL if is_evidence else BLUEFILL if is_other else WHITE,
                line=ORANGE if is_evidence else BLUE if is_other else MIDGRAY,
                crossed=not (is_evidence or is_other),
                name=f"After pruning r{row+1} c{col+1}",
            )

    c.rect(7.75, 3.25, 7.93, 3.43, fill=ORANGEFILL, line=ORANGE, width=0.45)
    c.text(8.00, 3.16, 9.90, 3.52, "selected evidence", size=9.5, align=PP_ALIGN.LEFT)
    c.rect(10.16, 3.25, 10.34, 3.43, fill=BLUEFILL, line=BLUE, width=0.45)
    c.text(10.41, 3.16, 12.24, 3.52, "other kept", size=10.0, align=PP_ALIGN.LEFT)
    add_token_cell(c, 12.26, 3.25, 0.18, 0.18, fill=WHITE, line=MIDGRAY, crossed=True, name="Pruned legend token")
    c.text(12.50, 3.16, 13.32, 3.52, "pruned", size=9.5, align=PP_ALIGN.LEFT)
    c.group("Stage 3 - budgeted token selection")


def build_stage_4(c: Canvas) -> None:
    c.clear_group()
    c.rect(13.85, 2.18, 17.82, 9.82, fill=WHITE, line=INK, width=1.0, rounded=True, name="Stage 4 panel")
    add_step_header(c, 4, 14.28, "Pruned MLLM Inference", (14.54, 9.26, 17.62, 9.66))
    c.text(14.12, 8.64, 17.56, 9.04, "Token sequence to LLM", size=13.0)

    c.line(14.20, 8.06, 15.70, 8.06, color=BLUE, width=0.8)
    c.line(14.20, 8.00, 14.20, 8.13, color=BLUE, width=0.8)
    c.line(15.70, 8.00, 15.70, 8.13, color=BLUE, width=0.8)
    c.text(14.20, 8.12, 15.70, 8.58, "visual prefix\n(shorter)", size=12.0, color=BLUE)
    c.line(15.94, 8.06, 17.47, 8.06, color=GREEN, width=0.8)
    c.line(15.94, 8.00, 15.94, 8.13, color=GREEN, width=0.8)
    c.line(17.47, 8.00, 17.47, 8.13, color=GREEN, width=0.8)
    c.text(15.94, 8.12, 17.47, 8.58, "text tokens\n(unchanged)", size=12.0, color=GREEN)

    visual_colors = [BLUEFILL, BLUEFILL, ORANGEFILL, ORANGEFILL, BLUEFILL]
    for idx, fill in enumerate(visual_colors):
        c.rect(14.25 + 0.27 * idx, 7.67, 14.43 + 0.27 * idx, 7.92, fill=fill, line=BLUE, width=0.45, name=f"Visual prefix token {idx+1}")
    c.text(15.55, 7.62, 15.98, 7.97, "...", size=12.0)
    for idx in range(5):
        c.rect(16.02 + 0.27 * idx, 7.67, 16.20 + 0.27 * idx, 7.92, fill=GREENFILL, line=GREEN, width=0.45, name=f"Text token {idx+1}")
    c.arrow_shape(15.75, 6.98, 15.93, 7.46, direction="down", fill=INK)

    c.labeled_rect(14.30, 5.26, 17.38, 6.93, "", fill=GREENFILL, line=GREEN, width=1.0, name="Editable LLM module")
    c.text(14.68, 6.28, 17.00, 6.76, "LLM", size=21.0, bold=True)
    c.text(14.68, 5.98, 17.00, 6.32, "(unchanged)", size=11.0)
    nodes = [(14.86, 5.69), (15.32, 5.51), (15.80, 5.82), (16.27, 5.49), (16.79, 5.78)]
    edges = [(0, 1), (1, 2), (2, 3), (3, 4), (0, 2), (2, 4), (1, 3)]
    for a, b in edges:
        c.line(*nodes[a], *nodes[b], color=GREEN, width=0.55)
    for idx, (cx, cy) in enumerate(nodes):
        c.circle(cx, cy, 0.16, fill=GREENFILL, line=GREEN, width=0.55, name=f"LLM node {idx+1}")
    c.arrow_shape(15.75, 4.82, 15.93, 5.24, direction="down", fill=INK)
    c.labeled_rect(14.53, 4.18, 17.15, 4.78, "Yes / No answer", fill=WHITE, line=MIDGRAY, width=0.8, size=15.0, bold=True, name="Editable model output")

    bullets = [
        ("shorter visual prefix", BLUE, 3.69),
        ("faster prefill", PURPLE, 3.30),
        ("lower memory", ORANGE, 2.91),
    ]
    for idx, (label, color, cy) in enumerate(bullets):
        tri = c.slide.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE, x(14.48), top_from_y(cy + 0.09), x(0.12), x(0.18))
        tri.fill.solid()
        tri.fill.fore_color.rgb = rgb(color)
        tri.line.fill.background()
        c.track(tri, f"Benefit bullet {idx+1}")
        c.text(14.64, cy - 0.16, 17.35, cy + 0.18, label, size=12.0, color=color, bold=True, align=PP_ALIGN.LEFT)
    c.group("Stage 4 - pruned MLLM inference")


def build_audit_band(c: Canvas) -> None:
    c.clear_group()
    c.text(7.15, 1.82, 10.85, 2.18, "Evidence Audit", size=17.0, bold=True, name="Audit title")
    c.rect(0.75, 0.18, 17.25, 1.66, fill=WHITE, line=INK, width=0.8, rounded=True, name="Audit band")
    audit_boxes = [
        (0.92, 2.86, "Accuracy / hFPR"),
        (3.93, 2.86, "Geometric\ncoverage\nPosECR / NegSRC"),
        (6.94, 2.86, "Interventional\ntests"),
        (9.95, 2.86, "BBox occlusion"),
        (12.96, 3.98, "Real CUDA speed"),
    ]
    for idx, (x0, width, label) in enumerate(audit_boxes):
        c.rect(x0, 0.38, x0 + width, 1.36, fill=LIGHTGRAY, line=MIDGRAY, width=0.55, rounded=True, name=f"Audit module {idx+1}")
        c.text(x0 + 0.70, 0.48, x0 + width - 0.08, 1.25, label, size=13.0, bold=True, align=PP_ALIGN.LEFT, name=f"Audit module {idx+1} label")

    c.circle(1.25, 0.87, 0.40, fill=WHITE, line=BLUE, width=0.9)
    c.line(1.25, 0.55, 1.25, 1.19, color=BLUE, width=0.9)
    c.line(0.93, 0.87, 1.57, 0.87, color=BLUE, width=0.9)

    c.circle(4.30, 0.87, 0.44, fill=GREENFILL, line=GREEN, width=0.9)
    c.line(4.30, 0.87, 4.30, 1.09, color=GREEN, width=0.8)
    c.line(4.30, 0.87, 4.49, 0.76, color=GREEN, width=0.8)

    c.line(7.28, 0.60, 7.28, 1.14, color=PURPLE, width=1.0)
    c.line(7.28, 0.60, 7.62, 0.84, color=PURPLE, width=1.0)
    c.line(7.62, 0.84, 7.94, 0.53, color=PURPLE, width=1.0)

    c.rect(10.30, 0.62, 10.79, 1.11, fill=WHITE, line=ORANGE, width=1.0, dashed=True, name="BBox audit icon")

    c.circle(13.58, 0.87, 0.48, fill=WHITE, line=BLUE, width=0.9)
    c.rect(13.25, 0.50, 13.91, 0.84, fill=LIGHTGRAY, line=LIGHTGRAY, width=0)
    c.line(13.58, 0.87, 13.76, 1.04, color=BLUE, width=1.0)
    c.group("Evidence audit band")


def add_stage_arrows(slide) -> None:
    c = Canvas(slide)
    for idx, (x1, x2) in enumerate(((3.57, 3.82), (7.14, 7.39), (13.57, 13.82)), start=1):
        c.arrow_shape(x1, 5.89, x2, 6.07, fill=INK, name=f"Stage transition {idx}")
    c.line(17.48, 4.48, 17.48, 1.90, color=INK, width=0.65, dashed=True)
    c.line(17.48, 1.90, 11.62, 1.90, color=INK, width=0.65, dashed=True)
    c.line(11.62, 1.90, 11.62, 2.16, color=INK, width=0.65, dashed=True, arrow=True, name="Audit feedback arrow")


def audit_pptx(path: Path) -> list[str]:
    defects: list[str] = []
    reopened = Presentation(path)
    if len(reopened.slides) != 1:
        defects.append(f"high: expected 1 slide, found {len(reopened.slides)}")
    slide = reopened.slides[0]
    if len(slide.shapes) < 8:
        defects.append(f"high: suspiciously few top-level shapes ({len(slide.shapes)})")
    for shape in slide.shapes:
        if shape.left < 0 or shape.top < 0:
            defects.append(f"high: shape outside top/left bounds: {shape.name}")
        if shape.left + shape.width > reopened.slide_width + 1:
            defects.append(f"high: shape outside right bound: {shape.name}")
        if shape.top + shape.height > reopened.slide_height + 1:
            defects.append(f"high: shape outside bottom bound: {shape.name}")
    group_names = {shape.name for shape in slide.shapes if shape.shape_type == 6}
    expected_groups = {
        "Stage 1 - OCR-critical query",
        "Stage 2 - projected visual tokens",
        "Stage 3 - budgeted token selection",
        "Stage 4 - pruned MLLM inference",
        "Evidence audit band",
    }
    missing = expected_groups - group_names
    if missing:
        defects.append(f"high: missing editable groups: {sorted(missing)}")
    for shape in iter_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
            continue
        run_sizes = [
            run.font.size.pt
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.font.size is not None
        ]
        if run_sizes and min(run_sizes) < 9.0:
            defects.append(
                f"high: text below 9 pt in editable PPTX: {shape.name} "
                f"({min(run_sizes):.1f} pt)"
            )
    return defects


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == 6:
            yield from iter_shapes(shape.shapes)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(WHITE)

    title_canvas = Canvas(slide)
    title_canvas.text(
        2.5,
        10.02,
        15.5,
        10.62,
        "Spatial-Provenance Audit for Visual Token Pruning",
        size=25.0,
        bold=True,
        name="Figure title",
    )

    canvas = Canvas(slide)
    build_stage_1(canvas)
    build_stage_2(canvas)
    build_stage_3(canvas)
    build_stage_4(canvas)
    build_audit_band(canvas)
    add_stage_arrows(slide)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    defects = audit_pptx(OUTPUT)
    if defects:
        raise RuntimeError("\n".join(defects))
    print(f"Wrote {OUTPUT}")
    return OUTPUT


if __name__ == "__main__":
    build()
